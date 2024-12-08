import os
import concurrent.futures
import docx  # For .docx
import openpyxl  # For .xlsx, .xlsm, etc.
import PyPDF2  # For .pdf
import pptx  # For .pptx
import time

def process_word_document(filepath, output_dir):
    """Extracts text from a Word document (.docx)."""
    try:
        doc = docx.Document(filepath)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        output_filename = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
        output_filepath = os.path.join(output_dir, output_filename)
        with open(output_filepath, "w", encoding="utf-8") as outfile:
            outfile.write(text)
    except Exception as e:
        print(f"Error processing Word document {filepath}: {e}")

def process_excel_document(filepath, output_dir):
    """Extracts text from an Excel document (.xlsx, .xlsm, etc.)."""
    try:
        workbook = openpyxl.load_workbook(filepath, read_only=True)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + "\t"  # Add tab as delimiter
                text += "\n"

        output_filename = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
        output_filepath = os.path.join(output_dir, output_filename)
        with open(output_filepath, "w", encoding="utf-8") as outfile:
            outfile.write(text)
    except Exception as e:
        print(f"Error processing Excel document {filepath}: {e}")

def process_pdf_document(filepath, output_dir):
    """Extracts text from a PDF document (.pdf)."""
    try:
        with open(filepath, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()

        output_filename = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
        output_filepath = os.path.join(output_dir, output_filename)
        with open(output_filepath, "w", encoding="utf-8") as outfile:
            outfile.write(text)
    except Exception as e:
        print(f"Error processing PDF document {filepath}: {e}")

def process_ppt_document(filepath, output_dir):
    """Extracts text from a PowerPoint document (.pptx)."""
    try:
        prs = pptx.Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        output_filename = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
        output_filepath = os.path.join(output_dir, output_filename)
        with open(output_filepath, "w", encoding="utf-8") as outfile:
            outfile.write(text)
    except Exception as e:
        print(f"Error processing PowerPoint document {filepath}: {e}")

def process_file(filepath, output_dir):
    """Processes a single file based on its extension."""
    if filepath.lower().endswith((".docx")):
        process_word_document(filepath, output_dir)
    elif filepath.lower().endswith((".xlsx", ".xlsm", ".xltx", ".xltm")):
        process_excel_document(filepath, output_dir)
    elif filepath.lower().endswith(".pdf"):
        process_pdf_document(filepath, output_dir)
    elif filepath.lower().endswith(".pptx"):
        process_ppt_document(filepath, output_dir)
    else:
        print(f"Unsupported file type: {filepath}")

def process_directory(input_dir, output_dir, max_workers=None):
    """Processes all files in a directory (and subdirectories) in parallel."""
    start_time = time.time()
    file_list = []
    for root, _, files in os.walk(input_dir):
        for file in files:
            filepath = os.path.join(root, file)
            file_list.append(filepath)

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)

        # Submit each file for processing
        futures = [executor.submit(process_file, filepath, output_dir) for filepath in file_list]

        # Wait for all tasks to complete
        for future in concurrent.futures.as_completed(futures):
            try:
                future.result()  # Get the result (or exception) of the task
            except Exception as e:
                print(f"Error in a worker thread: {e}")

    end_time = time.time()
    print(f"Total processing time: {end_time - start_time:.2f} seconds")

if __name__ == "__main__":
    input_directory = "./documents"  # Replace with your input directory
    output_directory = "./output"  # Replace with your desired output directory
    max_workers = 4  # Adjust the number of worker threads as needed (None uses the default)

    process_directory(input_directory, output_directory, max_workers)
